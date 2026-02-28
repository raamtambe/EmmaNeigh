#!/usr/bin/env python3
"""
EmmaNeigh - Table Comparison / Redline
v5.3.4: Compare tables across documents and generate redline markup

Supports extracting tables from: PDF, DOCX, PPTX, XLSX
Output: XLSX (color-coded comparison showing actual row/cell changes)

Core feature: Content-based row matching (not position-based).
When a row is inserted mid-table, only that insertion is marked — subsequent
rows are matched by content fingerprint, not by position. This solves the
"every row after an insert shows as changed" problem.

Table comparison approach:
1. Extract all tables from both documents
2. Match tables between documents (by header similarity + content overlap)
3. Match rows by content fingerprint (hash of cell values)
4. Fallback: match by key columns (first 1-2 columns, e.g. party name, date)
5. Cell-level diff for matched rows
6. Output: Excel workbook with color-coded changes
"""

import os
import sys
import json
import hashlib
import difflib
import re
from datetime import datetime
from typing import List, Dict, Tuple, Optional, Any
from dataclasses import dataclass, field
from concurrent.futures import ProcessPoolExecutor, as_completed
import multiprocessing as mp

# Document processing libraries
try:
    import fitz  # PyMuPDF
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from openpyxl import load_workbook
    from openpyxl import Workbook
    from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False


# =============================================================================
# Progress Emission
# =============================================================================

def emit(msg_type, **kwargs):
    """Output JSON message to stdout for the Electron app."""
    print(json.dumps({"type": msg_type, **kwargs}), flush=True)


# =============================================================================
# Text Normalization
# =============================================================================

def normalize_text(text: str) -> str:
    """Normalize text for comparison."""
    if not text:
        return ""
    text = re.sub(r'[\xa0\t\r]+', ' ', text)
    text = re.sub(r' +', ' ', text)
    text = text.replace('\u2018', "'").replace('\u2019', "'")
    text = text.replace('\u201c', '"').replace('\u201d', '"')
    text = text.replace('\u2013', '-').replace('\u2014', '-')
    return text.strip()


def texts_are_equivalent(text1: str, text2: str) -> bool:
    """Check if two texts are effectively the same after normalization."""
    return normalize_text(text1).lower() == normalize_text(text2).lower()


# =============================================================================
# Data Classes
# =============================================================================

@dataclass
class Cell:
    """Table cell with position and content."""
    row: int
    col: int
    text: str

    @property
    def fingerprint(self) -> str:
        return hashlib.md5(normalize_text(self.text).lower().encode()).hexdigest()[:8]


@dataclass
class TableRow:
    """Table row with cells."""
    index: int
    cells: List[Cell]
    is_header: bool = False

    @property
    def fingerprint(self) -> str:
        """Content-based fingerprint — position independent."""
        cell_fps = [c.fingerprint for c in self.cells]
        return hashlib.md5('|'.join(cell_fps).encode()).hexdigest()[:16]

    @property
    def text_content(self) -> str:
        """Concatenated text for fuzzy matching."""
        return ' | '.join(normalize_text(c.text) for c in self.cells)

    @property
    def key_values(self) -> List[str]:
        """First 3 columns — often contain unique identifiers."""
        return [normalize_text(c.text).lower() for c in self.cells[:3]]


@dataclass
class Table:
    """Table structure with rows."""
    id: str
    rows: List[TableRow]
    position: Tuple[int, int]  # (page/sheet, index_within_page)

    @property
    def header_row(self) -> Optional[TableRow]:
        if self.rows and self.rows[0].is_header:
            return self.rows[0]
        return None

    @property
    def header_fingerprint(self) -> str:
        if self.header_row:
            return self.header_row.fingerprint
        return ""

    @property
    def column_count(self) -> int:
        return max(len(row.cells) for row in self.rows) if self.rows else 0

    @property
    def data_rows(self) -> List[TableRow]:
        """All non-header rows."""
        return [r for r in self.rows if not r.is_header]


@dataclass
class CellChange:
    """Change in a table cell."""
    change_type: str  # 'unchanged', 'added', 'deleted', 'modified'
    original_value: Optional[str] = None
    modified_value: Optional[str] = None
    row: int = 0
    col: int = 0


@dataclass
class RowChange:
    """Change in a table row."""
    change_type: str  # 'unchanged', 'added', 'deleted', 'modified', 'moved'
    original_index: Optional[int] = None
    modified_index: Optional[int] = None
    cells: List[CellChange] = field(default_factory=list)
    moved_from: Optional[int] = None


@dataclass
class TableDiff:
    """Complete diff between two tables."""
    table_name: str
    row_changes: List[RowChange]
    column_mapping: Dict[int, int]
    is_resorted: bool
    added_columns: List[int]
    deleted_columns: List[int]
    stats: Dict[str, int] = field(default_factory=dict)


# =============================================================================
# Content Extractors — Tables Only
# =============================================================================

def extract_tables_from_pdf(file_path: str) -> List[Table]:
    """Extract tables from PDF using PyMuPDF."""
    if not HAS_FITZ:
        raise ImportError("PyMuPDF (fitz) not installed")

    tables = []
    doc = fitz.open(file_path)

    for page_num, page in enumerate(doc):
        try:
            page_tables = page.find_tables()
            for table_idx, table_data in enumerate(page_tables):
                rows = []
                extracted = table_data.extract()
                for row_idx, row_data in enumerate(extracted):
                    cells = [Cell(row=row_idx, col=col_idx, text=cell_text or "")
                             for col_idx, cell_text in enumerate(row_data)]
                    rows.append(TableRow(index=row_idx, cells=cells, is_header=(row_idx == 0)))
                if rows:
                    tables.append(Table(
                        id=f"Page{page_num + 1}_Table{table_idx + 1}",
                        rows=rows,
                        position=(page_num, table_idx)
                    ))
        except Exception:
            pass

    doc.close()
    return tables


def extract_tables_from_docx(file_path: str) -> List[Table]:
    """Extract tables from Word document."""
    if not HAS_DOCX:
        raise ImportError("python-docx not installed")

    tables = []
    doc = Document(file_path)

    for table_idx, docx_table in enumerate(doc.tables):
        rows = []
        for row_idx, docx_row in enumerate(docx_table.rows):
            cells = [Cell(row=row_idx, col=col_idx, text=docx_cell.text)
                     for col_idx, docx_cell in enumerate(docx_row.cells)]
            rows.append(TableRow(index=row_idx, cells=cells, is_header=(row_idx == 0)))
        if rows:
            tables.append(Table(
                id=f"Table{table_idx + 1}",
                rows=rows,
                position=(0, table_idx)
            ))

    return tables


def extract_tables_from_pptx(file_path: str) -> List[Table]:
    """Extract tables from PowerPoint."""
    if not HAS_PPTX:
        raise ImportError("python-pptx not installed")

    tables = []
    prs = Presentation(file_path)

    for slide_idx, slide in enumerate(prs.slides):
        table_count = 0
        for shape in slide.shapes:
            if shape.has_table:
                pptx_table = shape.table
                rows = []
                for row_idx in range(len(pptx_table.rows)):
                    cells = [Cell(row=row_idx, col=col_idx, text=pptx_table.cell(row_idx, col_idx).text)
                             for col_idx in range(len(pptx_table.columns))]
                    rows.append(TableRow(index=row_idx, cells=cells, is_header=(row_idx == 0)))
                if rows:
                    table_count += 1
                    tables.append(Table(
                        id=f"Slide{slide_idx + 1}_Table{table_count}",
                        rows=rows,
                        position=(slide_idx, table_count - 1)
                    ))

    return tables


def extract_tables_from_xlsx(file_path: str) -> List[Table]:
    """Extract tables from Excel. Each sheet = one table."""
    if not HAS_OPENPYXL:
        raise ImportError("openpyxl not installed")

    tables = []
    wb = load_workbook(file_path, data_only=True)

    for sheet_idx, sheet_name in enumerate(wb.sheetnames):
        sheet = wb[sheet_name]
        rows = []
        for row_idx, row in enumerate(sheet.iter_rows()):
            cells = []
            has_content = False
            for col_idx, cell in enumerate(row):
                cell_value = str(cell.value) if cell.value is not None else ""
                if cell_value:
                    has_content = True
                cells.append(Cell(row=row_idx, col=col_idx, text=cell_value))
            if has_content:
                rows.append(TableRow(index=row_idx, cells=cells, is_header=(row_idx == 0)))
        if rows:
            tables.append(Table(
                id=sheet_name,
                rows=rows,
                position=(sheet_idx, 0)
            ))

    wb.close()
    return tables


def extract_tables(file_path: str) -> List[Table]:
    """Extract all tables from a document."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.pdf':
        return extract_tables_from_pdf(file_path)
    elif ext in ['.docx', '.doc']:
        return extract_tables_from_docx(file_path)
    elif ext in ['.pptx', '.ppt']:
        return extract_tables_from_pptx(file_path)
    elif ext in ['.xlsx', '.xls']:
        return extract_tables_from_xlsx(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")


# =============================================================================
# Table Matching (between documents)
# =============================================================================

def match_tables(orig_tables: List[Table], mod_tables: List[Table]) -> List[Tuple[Optional[Table], Optional[Table]]]:
    """Match original tables to modified tables by content similarity."""
    matches = []
    used_modified = set()

    for orig in orig_tables:
        best_match = None
        best_score = 0.3  # Low threshold — prefer matching over orphaning

        for i, mod in enumerate(mod_tables):
            if i in used_modified:
                continue
            score = compute_table_similarity(orig, mod)
            if score > best_score:
                best_score = score
                best_match = (i, mod)

        if best_match:
            used_modified.add(best_match[0])
            matches.append((orig, best_match[1]))
        else:
            matches.append((orig, None))

    for i, mod in enumerate(mod_tables):
        if i not in used_modified:
            matches.append((None, mod))

    return matches


def compute_table_similarity(t1: Table, t2: Table) -> float:
    """Compute similarity between two tables."""
    scores = []

    # Header similarity (50% weight) — strongest signal
    if t1.header_fingerprint and t2.header_fingerprint:
        if t1.header_fingerprint == t2.header_fingerprint:
            scores.append(1.0 * 0.5)
        else:
            # Partial header match
            h1_cells = [normalize_text(c.text).lower() for c in (t1.header_row.cells if t1.header_row else [])]
            h2_cells = [normalize_text(c.text).lower() for c in (t2.header_row.cells if t2.header_row else [])]
            if h1_cells and h2_cells:
                overlap = len(set(h1_cells) & set(h2_cells))
                scores.append((overlap / max(len(h1_cells), len(h2_cells))) * 0.5)

    # Content overlap (30% weight)
    fp1 = set(row.fingerprint for row in t1.data_rows)
    fp2 = set(row.fingerprint for row in t2.data_rows)
    if fp1 or fp2:
        intersection = len(fp1 & fp2)
        union = len(fp1 | fp2)
        scores.append((intersection / union if union > 0 else 0) * 0.3)

    # Position (20% weight)
    if t1.position[0] == t2.position[0]:
        scores.append(0.2)

    return sum(scores)


# =============================================================================
# Row Matching — Content-Based (solves the "inserted row cascades" problem)
# =============================================================================

def match_rows_by_content(orig_table: Table, mod_table: Table, col_mapping: Dict[int, int]) -> Dict[int, Tuple[int, float]]:
    """
    Match rows between tables by CONTENT, not position.

    This is the key algorithm that solves the "one inserted row makes
    everything after it show as changed" problem.

    Strategy:
    1. Exact fingerprint match (hash of all cell values) — O(n)
    2. Key column match (first 1-3 columns) — catches rows where only some cells changed
    3. Fuzzy text similarity — last resort for heavily modified rows

    Returns: {orig_row_index: (mod_row_index, confidence)}
    """
    matches = {}
    used_modified = set()

    orig_data = [(i, r) for i, r in enumerate(orig_table.rows) if not r.is_header]
    mod_data = [(i, r) for i, r in enumerate(mod_table.rows) if not r.is_header]

    # Build lookup indices for speed
    mod_fp_index = {}  # fingerprint → list of (mod_idx, mod_row)
    for mod_idx, mod_row in mod_data:
        fp = mod_row.fingerprint
        mod_fp_index.setdefault(fp, []).append((mod_idx, mod_row))

    mod_key_index = {}  # (key1, key2) → list of (mod_idx, mod_row)
    for mod_idx, mod_row in mod_data:
        key = tuple(mod_row.key_values)
        mod_key_index.setdefault(key, []).append((mod_idx, mod_row))

    # === Pass 1: Exact fingerprint match ===
    for orig_idx, orig_row in orig_data:
        fp = orig_row.fingerprint
        if fp in mod_fp_index:
            for mod_idx, mod_row in mod_fp_index[fp]:
                if mod_idx not in used_modified:
                    matches[orig_idx] = (mod_idx, 1.0)
                    used_modified.add(mod_idx)
                    break

    # === Pass 2: Key column match ===
    for orig_idx, orig_row in orig_data:
        if orig_idx in matches:
            continue

        orig_key = tuple(orig_row.key_values)

        # Skip empty keys
        if all(k == '' for k in orig_key):
            continue

        if orig_key in mod_key_index:
            for mod_idx, mod_row in mod_key_index[orig_key]:
                if mod_idx not in used_modified:
                    matches[orig_idx] = (mod_idx, 0.8)
                    used_modified.add(mod_idx)
                    break

    # === Pass 3: Fuzzy similarity for remaining unmatched rows ===
    unmatched_orig = [(i, r) for i, r in orig_data if i not in matches]
    unmatched_mod = [(i, r) for i, r in mod_data if i not in used_modified]

    for orig_idx, orig_row in unmatched_orig:
        best_score = 0.6  # High threshold for fuzzy — avoid false matches
        best_mod_idx = None

        orig_text = orig_row.text_content

        for mod_idx, mod_row in unmatched_mod:
            if mod_idx in used_modified:
                continue

            mod_text = mod_row.text_content
            sim = difflib.SequenceMatcher(None, orig_text, mod_text).ratio()

            if sim > best_score:
                best_score = sim
                best_mod_idx = mod_idx

        if best_mod_idx is not None:
            matches[orig_idx] = (best_mod_idx, best_score)
            used_modified.add(best_mod_idx)

    return matches


# =============================================================================
# Column Matching
# =============================================================================

def detect_column_reorder(orig_table: Table, mod_table: Table) -> Dict[int, int]:
    """Detect column reordering by matching header cells."""
    if not orig_table.header_row or not mod_table.header_row:
        # No headers — assume identity mapping up to min column count
        n = min(orig_table.column_count, mod_table.column_count)
        return {i: i for i in range(n)}

    orig_headers = [normalize_text(c.text).lower() for c in orig_table.header_row.cells]
    mod_headers = [normalize_text(c.text).lower() for c in mod_table.header_row.cells]

    mapping = {}
    used_mod = set()

    # Exact match first
    for i, oh in enumerate(orig_headers):
        if not oh:
            continue
        for j, mh in enumerate(mod_headers):
            if j in used_mod:
                continue
            if oh == mh:
                mapping[i] = j
                used_mod.add(j)
                break

    # Fuzzy match for remaining
    for i, oh in enumerate(orig_headers):
        if i in mapping or not oh:
            continue
        best_j = None
        best_sim = 0.7
        for j, mh in enumerate(mod_headers):
            if j in used_mod or not mh:
                continue
            sim = difflib.SequenceMatcher(None, oh, mh).ratio()
            if sim > best_sim:
                best_sim = sim
                best_j = j
        if best_j is not None:
            mapping[i] = best_j
            used_mod.add(best_j)

    # Remaining unmapped columns: identity if within range
    for i in range(orig_table.column_count):
        if i not in mapping and i < mod_table.column_count and i not in used_mod:
            mapping[i] = i
            used_mod.add(i)

    return mapping


# =============================================================================
# Table Diff Generation
# =============================================================================

def diff_tables(orig: Table, mod: Table) -> TableDiff:
    """Generate detailed diff between two matched tables."""
    col_mapping = detect_column_reorder(orig, mod)

    orig_cols = set(range(orig.column_count))
    mod_cols = set(range(mod.column_count))
    mapped_orig = set(col_mapping.keys())
    mapped_mod = set(col_mapping.values())

    deleted_columns = sorted(orig_cols - mapped_orig)
    added_columns = sorted(mod_cols - mapped_mod)

    # Content-based row matching
    row_matches = match_rows_by_content(orig, mod, col_mapping)

    # Detect resorting
    matched_orig_indices = sorted(row_matches.keys())
    matched_mod_indices = [row_matches[i][0] for i in matched_orig_indices]
    is_resorted = matched_mod_indices != sorted(matched_mod_indices) if matched_mod_indices else False

    # Generate row changes
    row_changes = []
    used_mod_rows = set()

    # Process original rows in order
    for orig_idx, orig_row in enumerate(orig.rows):
        if orig_row.is_header:
            # Header row — compare but don't flag as change
            if mod.header_row:
                cell_changes = []
                for oc, mc in col_mapping.items():
                    if oc < len(orig_row.cells):
                        ov = orig_row.cells[oc].text
                        mv = mod.header_row.cells[mc].text if mc < len(mod.header_row.cells) else ""
                        ct = 'unchanged' if texts_are_equivalent(ov, mv) else 'modified'
                        cell_changes.append(CellChange(change_type=ct, original_value=ov, modified_value=mv, row=orig_idx, col=oc))
                row_changes.append(RowChange(change_type='unchanged', original_index=orig_idx, modified_index=0, cells=cell_changes))
                used_mod_rows.add(0)
            continue

        if orig_idx in row_matches:
            mod_idx, confidence = row_matches[orig_idx]
            mod_row = mod.rows[mod_idx]
            used_mod_rows.add(mod_idx)

            # Cell-level comparison
            cell_changes = []
            has_cell_change = False

            for orig_col, mod_col in col_mapping.items():
                if orig_col >= len(orig_row.cells):
                    continue
                orig_cell = orig_row.cells[orig_col]

                if mod_col < len(mod_row.cells):
                    mod_cell = mod_row.cells[mod_col]
                    if texts_are_equivalent(orig_cell.text, mod_cell.text):
                        ct = 'unchanged'
                    else:
                        ct = 'modified'
                        has_cell_change = True
                    cell_changes.append(CellChange(
                        change_type=ct,
                        original_value=orig_cell.text,
                        modified_value=mod_cell.text,
                        row=orig_idx, col=orig_col
                    ))
                else:
                    cell_changes.append(CellChange(
                        change_type='deleted',
                        original_value=orig_cell.text,
                        row=orig_idx, col=orig_col
                    ))
                    has_cell_change = True

            # Determine row change type
            if mod_idx != orig_idx and has_cell_change:
                row_type = 'modified'  # Moved AND modified
            elif mod_idx != orig_idx:
                row_type = 'moved'
            elif has_cell_change:
                row_type = 'modified'
            else:
                row_type = 'unchanged'

            row_changes.append(RowChange(
                change_type=row_type,
                original_index=orig_idx,
                modified_index=mod_idx,
                cells=cell_changes,
                moved_from=orig_idx if mod_idx != orig_idx else None
            ))
        else:
            # Deleted row
            cell_changes = [CellChange(change_type='deleted', original_value=c.text, row=orig_idx, col=c.col)
                            for c in orig_row.cells]
            row_changes.append(RowChange(change_type='deleted', original_index=orig_idx, cells=cell_changes))

    # Added rows (in modified but not matched)
    for mod_idx, mod_row in enumerate(mod.rows):
        if mod_idx not in used_mod_rows and not mod_row.is_header:
            cell_changes = [CellChange(change_type='added', modified_value=c.text, row=mod_idx, col=c.col)
                            for c in mod_row.cells]
            row_changes.append(RowChange(change_type='added', modified_index=mod_idx, cells=cell_changes))

    # Compute stats
    stats = {
        'total_rows_orig': len(orig.data_rows),
        'total_rows_mod': len(mod.data_rows),
        'unchanged': sum(1 for r in row_changes if r.change_type == 'unchanged'),
        'added': sum(1 for r in row_changes if r.change_type == 'added'),
        'deleted': sum(1 for r in row_changes if r.change_type == 'deleted'),
        'modified': sum(1 for r in row_changes if r.change_type == 'modified'),
        'moved': sum(1 for r in row_changes if r.change_type == 'moved'),
    }

    return TableDiff(
        table_name=orig.id,
        row_changes=row_changes,
        column_mapping=col_mapping,
        is_resorted=is_resorted,
        added_columns=added_columns,
        deleted_columns=deleted_columns,
        stats=stats
    )


def compare_all_tables(orig_tables: List[Table], mod_tables: List[Table]) -> List[Tuple[Optional[Table], Optional[Table], Optional[TableDiff]]]:
    """Compare all tables between two documents."""
    results = []
    table_pairs = match_tables(orig_tables, mod_tables)

    for orig, mod in table_pairs:
        if orig is None:
            results.append((None, mod, None))
        elif mod is None:
            results.append((orig, None, None))
        else:
            td = diff_tables(orig, mod)
            results.append((orig, mod, td))

    return results


# =============================================================================
# Excel Output — Color-Coded Comparison
# =============================================================================

FILL_ADDED = PatternFill(start_color="CCE5FF", end_color="CCE5FF", fill_type="solid")
FILL_DELETED = PatternFill(start_color="FFCCCC", end_color="FFCCCC", fill_type="solid")
FILL_MODIFIED = PatternFill(start_color="FFFFCC", end_color="FFFFCC", fill_type="solid")
FILL_MOVED = PatternFill(start_color="CCFFCC", end_color="CCFFCC", fill_type="solid")
FILL_HEADER = PatternFill(start_color="D9D9D9", end_color="D9D9D9", fill_type="solid")

FONT_ADDED = Font(color="0000CC")
FONT_DELETED = Font(color="CC0000", strikethrough=True)
FONT_MODIFIED_OLD = Font(color="CC0000", strikethrough=True, size=9)
FONT_MODIFIED_NEW = Font(color="0000CC", size=9)
FONT_HEADER = Font(bold=True)
FONT_LABEL = Font(bold=True, size=9, color="555555")

THIN_BORDER = Border(
    left=Side(style='thin', color='AAAAAA'),
    right=Side(style='thin', color='AAAAAA'),
    top=Side(style='thin', color='AAAAAA'),
    bottom=Side(style='thin', color='AAAAAA')
)


def generate_output_xlsx(table_results: List[Tuple],
                         output_path: str,
                         orig_file: str,
                         mod_file: str):
    """Generate Excel workbook with color-coded table comparison."""
    if not HAS_OPENPYXL:
        raise ImportError("openpyxl not installed")

    wb = Workbook()
    if 'Sheet' in wb.sheetnames:
        del wb['Sheet']

    # === Summary Sheet ===
    ws_sum = wb.create_sheet(title="Summary", index=0)
    ws_sum.cell(row=1, column=1, value="Table Comparison Report").font = Font(bold=True, size=14)
    ws_sum.cell(row=3, column=1, value="Original:").font = Font(bold=True)
    ws_sum.cell(row=3, column=2, value=os.path.basename(orig_file))
    ws_sum.cell(row=4, column=1, value="Modified:").font = Font(bold=True)
    ws_sum.cell(row=4, column=2, value=os.path.basename(mod_file))
    ws_sum.cell(row=5, column=1, value="Generated:").font = Font(bold=True)
    ws_sum.cell(row=5, column=2, value=datetime.now().strftime('%Y-%m-%d %H:%M:%S'))

    # Legend
    row = 7
    ws_sum.cell(row=row, column=1, value="Color Legend:").font = Font(bold=True)
    row += 1
    for label, fill, font in [
        ("Row Added", FILL_ADDED, FONT_ADDED),
        ("Row Deleted", FILL_DELETED, FONT_DELETED),
        ("Cell Modified", FILL_MODIFIED, Font(color="CC6600")),
        ("Row Moved", FILL_MOVED, Font(color="006600")),
        ("Unchanged", PatternFill(), Font()),
    ]:
        c = ws_sum.cell(row=row, column=1, value=label)
        c.fill = fill
        c.font = font
        row += 1

    # Stats table
    row += 1
    headers = ["Table", "Orig Rows", "Mod Rows", "Unchanged", "Added", "Deleted", "Modified", "Moved"]
    for ci, h in enumerate(headers, 1):
        c = ws_sum.cell(row=row, column=ci, value=h)
        c.font = Font(bold=True)
        c.fill = FILL_HEADER
    row += 1

    for orig_table, mod_table, table_diff in table_results:
        name = (orig_table or mod_table).id
        ws_sum.cell(row=row, column=1, value=name)

        if table_diff:
            s = table_diff.stats
            ws_sum.cell(row=row, column=2, value=s.get('total_rows_orig', 0))
            ws_sum.cell(row=row, column=3, value=s.get('total_rows_mod', 0))
            ws_sum.cell(row=row, column=4, value=s.get('unchanged', 0))
            ws_sum.cell(row=row, column=5, value=s.get('added', 0))
            ws_sum.cell(row=row, column=6, value=s.get('deleted', 0))
            ws_sum.cell(row=row, column=7, value=s.get('modified', 0))
            ws_sum.cell(row=row, column=8, value=s.get('moved', 0))
        elif orig_table is None:
            ws_sum.cell(row=row, column=5, value="Entire table added")
        elif mod_table is None:
            ws_sum.cell(row=row, column=6, value="Entire table deleted")
        row += 1

    for col_letter in ['A', 'B', 'C', 'D', 'E', 'F', 'G', 'H']:
        ws_sum.column_dimensions[col_letter].width = 15

    # === Per-Table Sheets ===
    for orig_table, mod_table, table_diff in table_results:
        tbl = orig_table or mod_table
        sheet_name = tbl.id[:31]  # Excel max 31 chars

        # Ensure unique sheet name
        base_name = sheet_name
        counter = 1
        while sheet_name in wb.sheetnames:
            sheet_name = f"{base_name[:28]}_{counter}"
            counter += 1

        ws = wb.create_sheet(title=sheet_name)

        if orig_table is None and mod_table is not None:
            # Entire table added
            ws.cell(row=1, column=1, value="[ENTIRE TABLE ADDED]").font = FONT_ADDED
            for ri, row_obj in enumerate(mod_table.rows, 2):
                for cell in row_obj.cells:
                    c = ws.cell(row=ri, column=cell.col + 1, value=cell.text)
                    c.fill = FILL_ADDED
                    c.font = FONT_ADDED
                    c.border = THIN_BORDER
            _auto_width(ws)
            continue

        if mod_table is None and orig_table is not None:
            # Entire table deleted
            ws.cell(row=1, column=1, value="[ENTIRE TABLE DELETED]").font = FONT_DELETED
            for ri, row_obj in enumerate(orig_table.rows, 2):
                for cell in row_obj.cells:
                    c = ws.cell(row=ri, column=cell.col + 1, value=cell.text)
                    c.fill = FILL_DELETED
                    c.font = FONT_DELETED
                    c.border = THIN_BORDER
            _auto_width(ws)
            continue

        if table_diff is None:
            continue

        # Add a "Change" label column
        max_col = max(orig_table.column_count, mod_table.column_count)
        change_col = max_col + 2  # Leave a gap

        # Write header for change label
        ws.cell(row=1, column=change_col, value="Change").font = FONT_LABEL

        current_row = 1

        # Sort row_changes: header first, then by position
        sorted_changes = sorted(table_diff.row_changes,
                                key=lambda rc: (
                                    0 if rc.original_index == 0 or rc.modified_index == 0 else 1,
                                    rc.original_index if rc.original_index is not None else 99999,
                                    rc.modified_index if rc.modified_index is not None else 99999
                                ))

        for rc in sorted_changes:
            if rc.change_type == 'unchanged':
                for cc in rc.cells:
                    val = cc.modified_value if cc.modified_value is not None else cc.original_value or ""
                    c = ws.cell(row=current_row, column=cc.col + 1, value=val)
                    c.border = THIN_BORDER
                    if rc.original_index == 0:  # Header
                        c.fill = FILL_HEADER
                        c.font = FONT_HEADER

            elif rc.change_type == 'added':
                for cc in rc.cells:
                    c = ws.cell(row=current_row, column=cc.col + 1, value=cc.modified_value or "")
                    c.fill = FILL_ADDED
                    c.font = FONT_ADDED
                    c.border = THIN_BORDER
                ws.cell(row=current_row, column=change_col, value="ADDED").font = FONT_ADDED

            elif rc.change_type == 'deleted':
                for cc in rc.cells:
                    c = ws.cell(row=current_row, column=cc.col + 1, value=cc.original_value or "")
                    c.fill = FILL_DELETED
                    c.font = FONT_DELETED
                    c.border = THIN_BORDER
                ws.cell(row=current_row, column=change_col, value="DELETED").font = FONT_DELETED

            elif rc.change_type == 'modified':
                for cc in rc.cells:
                    if cc.change_type == 'modified':
                        # Show "old → new" in the cell
                        val = f"{cc.original_value} → {cc.modified_value}"
                        c = ws.cell(row=current_row, column=cc.col + 1, value=val)
                        c.fill = FILL_MODIFIED
                        c.font = Font(color="CC6600")
                    else:
                        val = cc.modified_value if cc.modified_value is not None else cc.original_value or ""
                        c = ws.cell(row=current_row, column=cc.col + 1, value=val)
                    c.border = THIN_BORDER
                ws.cell(row=current_row, column=change_col, value="MODIFIED").font = Font(color="CC6600")

            elif rc.change_type == 'moved':
                for cc in rc.cells:
                    val = cc.modified_value if cc.modified_value is not None else cc.original_value or ""
                    c = ws.cell(row=current_row, column=cc.col + 1, value=val)
                    c.fill = FILL_MOVED
                    c.border = THIN_BORDER
                label = f"MOVED (was row {rc.moved_from + 1})" if rc.moved_from is not None else "MOVED"
                ws.cell(row=current_row, column=change_col, value=label).font = Font(color="006600")

            current_row += 1

        _auto_width(ws)

    wb.save(output_path)


def _auto_width(ws):
    """Auto-fit column widths."""
    for col_cells in ws.columns:
        max_len = 0
        for cell in col_cells:
            if cell.value:
                max_len = max(max_len, len(str(cell.value)))
        ws.column_dimensions[col_cells[0].column_letter].width = min(max_len + 2, 60)


# =============================================================================
# Main Comparison Function
# =============================================================================

def compare_documents(original_path: str, modified_path: str, output_path: str) -> Dict:
    """Compare tables in two documents and generate Excel output."""

    emit("progress", percent=10, message=f"Extracting tables from {os.path.basename(original_path)}...")
    orig_tables = extract_tables(original_path)

    emit("progress", percent=30, message=f"Extracting tables from {os.path.basename(modified_path)}...")
    mod_tables = extract_tables(modified_path)

    if not orig_tables and not mod_tables:
        emit("progress", percent=100, message="No tables found in either document.")
        return {
            'output_path': output_path,
            'tables_compared': 0,
            'table_changes': 0,
            'message': 'No tables found in either document'
        }

    emit("progress", percent=60, message=f"Comparing {len(orig_tables)} original table(s) with {len(mod_tables)} modified table(s)...")
    table_results = compare_all_tables(orig_tables, mod_tables)

    emit("progress", percent=85, message="Generating comparison spreadsheet...")

    # Always output as xlsx
    if not output_path.endswith('.xlsx'):
        output_path = os.path.splitext(output_path)[0] + '.xlsx'

    generate_output_xlsx(table_results, output_path, original_path, modified_path)

    # Stats
    total_added = 0
    total_deleted = 0
    total_modified = 0
    for _, _, td in table_results:
        if td:
            total_added += td.stats.get('added', 0)
            total_deleted += td.stats.get('deleted', 0)
            total_modified += td.stats.get('modified', 0)

    return {
        'output_path': output_path,
        'tables_compared': len(orig_tables),
        'table_changes': sum(1 for _, _, d in table_results if d is not None),
        'rows_added': total_added,
        'rows_deleted': total_deleted,
        'rows_modified': total_modified
    }


# =============================================================================
# Batch Processing
# =============================================================================

def process_single_pair(args: Tuple[str, str, str]) -> Dict:
    original, modified, output = args
    try:
        result = compare_documents(original, modified, output)
        result['success'] = True
        return result
    except Exception as e:
        return {
            'success': False,
            'error': str(e),
            'original': original,
            'modified': modified
        }


def process_batch(pairs: List[Dict], output_folder: str) -> List[Dict]:
    results = []
    total = len(pairs)

    args_list = []
    for pair in pairs:
        original = pair['original']
        modified = pair['modified']
        orig_name = os.path.splitext(os.path.basename(original))[0]
        mod_name = os.path.splitext(os.path.basename(modified))[0]
        output_name = f"Redline_{orig_name}_vs_{mod_name}.xlsx"
        output_path = os.path.join(output_folder, output_name)
        args_list.append((original, modified, output_path))

    max_workers = max(1, mp.cpu_count() // 2)

    with ProcessPoolExecutor(max_workers=max_workers) as executor:
        futures = {executor.submit(process_single_pair, args): i
                   for i, args in enumerate(args_list)}

        completed = 0
        for future in as_completed(futures):
            idx = futures[future]
            try:
                result = future.result()
                results.append(result)
            except Exception as e:
                results.append({'success': False, 'error': str(e), 'index': idx})

            completed += 1
            emit("progress",
                 percent=int(completed / total * 100),
                 message=f"Completed {completed}/{total} comparisons")

    return results


# =============================================================================
# Main Entry Point
# =============================================================================

def main():
    if len(sys.argv) < 2:
        emit("error", message="Usage: document_redline.py <config_json_path>")
        sys.exit(1)

    config_path = sys.argv[1]
    if not os.path.isfile(config_path):
        emit("error", message=f"Config file not found: {config_path}")
        sys.exit(1)

    try:
        with open(config_path, 'r') as f:
            config = json.load(f)
    except json.JSONDecodeError as e:
        emit("error", message=f"Invalid JSON config: {str(e)}")
        sys.exit(1)

    try:
        if config.get('batch'):
            pairs = config.get('pairs', [])
            output_folder = config.get('output_folder', os.path.dirname(config_path))
            os.makedirs(output_folder, exist_ok=True)
            results = process_batch(pairs, output_folder)
            successful = sum(1 for r in results if r.get('success'))
            emit("result", success=True, mode="batch", total=len(pairs), successful=successful, results=results)
        else:
            original = config.get('original')
            modified = config.get('modified')
            output = config.get('output')

            if not original or not modified:
                emit("error", message="Missing 'original' or 'modified' path")
                sys.exit(1)

            if not output:
                orig_name = os.path.splitext(os.path.basename(original))[0]
                mod_name = os.path.splitext(os.path.basename(modified))[0]
                output = os.path.join(
                    os.path.dirname(original),
                    f"Redline_{orig_name}_vs_{mod_name}.xlsx"
                )

            result = compare_documents(original, modified, output)
            emit("result", success=True, mode="single", **result)

    except Exception as e:
        emit("error", message=f"Processing failed: {str(e)}")
        sys.exit(1)


if __name__ == "__main__":
    main()
