from pathlib import Path
import sys


TMP_VENDOR = Path("/tmp/emmaneigh-pptx")
if TMP_VENDOR.exists():
    sys.path.insert(0, str(TMP_VENDOR))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont


ROOT = Path(__file__).resolve().parent
DOCS = ROOT.parent
LOGOS = ROOT / "logos"

PX_PER_IN = 144
FONT_BODY = "Arial"
FONT_TITLE = "Arial Narrow"
FONT_BODY_FILE = "/System/Library/Fonts/Supplemental/Arial.ttf"
FONT_TITLE_FILE = "/System/Library/Fonts/Supplemental/Arial Narrow Bold.ttf"
_MEASURE_IMG = Image.new("RGB", (2400, 1600), "white")
_MEASURE_DRAW = ImageDraw.Draw(_MEASURE_IMG)


def px(value):
    return Inches(value / PX_PER_IN)


def rgb(hex_value):
    hex_value = hex_value.lstrip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def resolve_font(font_name):
    if "Avenir Next Condensed" in font_name:
        return FONT_TITLE, FONT_TITLE_FILE
    if "Avenir" in font_name:
        return FONT_BODY, FONT_BODY_FILE
    if font_name == FONT_TITLE:
        return FONT_TITLE, FONT_TITLE_FILE
    return FONT_BODY, FONT_BODY_FILE


def estimate_text_height(text, font_name, font_size, width_px, spacing=1.0):
    _, font_file = resolve_font(font_name)
    font = ImageFont.truetype(font_file, int(font_size))
    lines = []
    for paragraph in text.split("\n"):
        words = paragraph.split()
        if not words:
            lines.append("")
            continue
        current = words[0]
        for word in words[1:]:
            candidate = f"{current} {word}"
            if _MEASURE_DRAW.textlength(candidate, font=font) <= width_px:
                current = candidate
            else:
                lines.append(current)
                current = word
        lines.append(current)
    ascent, descent = font.getmetrics()
    line_height = ascent + descent
    spacing_px = max(0, int((spacing - 1.0) * line_height))
    return max(1, len(lines)) * line_height + max(0, len(lines) - 1) * spacing_px


def remove_placeholders(slide):
    for shape in list(slide.placeholders):
        sp = shape._element
        sp.getparent().remove(sp)


def add_rect(slide, x, y, w, h, fill, line=None, line_width=1.5, rounded=True, transparency=0.0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, px(x), px(y), px(w), px(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if transparency:
        shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color, width=2):
    line = slide.shapes.add_connector(1, px(x1), px(y1), px(x2), px(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_circle(slide, x, y, size, fill, line=None):
    return add_rect(slide, x, y, size, size, fill=fill, line=line, rounded=True)


def set_text_frame(tf, text, font_name, font_size, color, bold=False, align=PP_ALIGN.LEFT, valign=MSO_VERTICAL_ANCHOR.TOP, margins=(0, 0, 0, 0), spacing=1.1, wrap=True, autofit=False):
    if "Avenir Next Condensed" in font_name:
        actual_font = FONT_TITLE
        font_file = FONT_TITLE_FILE
    elif "Avenir" in font_name:
        actual_font = FONT_BODY
        font_file = FONT_BODY_FILE
    else:
        actual_font = font_name
        font_file = FONT_BODY_FILE if font_name == FONT_BODY else FONT_TITLE_FILE if font_name == FONT_TITLE else None
    tf.clear()
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = valign
    tf.margin_left = Pt(margins[0])
    tf.margin_right = Pt(margins[1])
    tf.margin_top = Pt(margins[2])
    tf.margin_bottom = Pt(margins[3])
    paragraphs = text.split("\n")
    first = True
    for para_text in paragraphs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = para_text
        p.alignment = align
        p.line_spacing = spacing
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = para_text
        run.font.name = actual_font
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    try:
        if autofit:
            tf.fit_text(font_family=actual_font, max_size=int(font_size), bold=bold, font_file=font_file)
    except Exception:
        pass


def add_textbox(slide, x, y, w, h, text, font_name=FONT_BODY, font_size=20, color="#ffffff", bold=False, align=PP_ALIGN.LEFT, valign=MSO_VERTICAL_ANCHOR.TOP, margins=(2, 2, 0, 0), spacing=1.0, wrap=True, autofit=False):
    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    set_text_frame(box.text_frame, text, font_name, font_size, color, bold=bold, align=align, valign=valign, margins=margins, spacing=spacing, wrap=wrap, autofit=autofit)
    return box


def add_picture_contain(slide, image_path, x, y, w, h):
    with Image.open(image_path) as im:
        img_w, img_h = im.size
    ratio = min(w / img_w, h / img_h)
    draw_w = img_w * ratio
    draw_h = img_h * ratio
    left = x + (w - draw_w) / 2
    top = y + (h - draw_h) / 2
    return slide.shapes.add_picture(str(image_path), px(left), px(top), width=px(draw_w), height=px(draw_h))


def add_logo_chip(slide, label, logo_name, x, y, w, h, fill, line, text_color, font_size=17):
    add_rect(slide, x, y, w, h, fill=fill, line=line, line_width=1.25)
    logo_box_w = min(40, h - 12)
    add_picture_contain(slide, LOGOS / logo_name, x + 12, y + 6, logo_box_w, h - 12)
    add_textbox(slide, x + 12 + logo_box_w + 10, y + 5, w - logo_box_w - 28, h - 10, label, font_size=font_size, color=text_color, bold=True, valign=MSO_VERTICAL_ANCHOR.MIDDLE, wrap=False, margins=(2, 2, 1, 1))


def add_text_chip(slide, label, x, y, w, h, fill, line, text_color, font_size=17):
    add_rect(slide, x, y, w, h, fill=fill, line=line, line_width=1.25)
    add_textbox(slide, x + 14, y + 5, w - 28, h - 10, label, font_size=font_size, color=text_color, bold=True, valign=MSO_VERTICAL_ANCHOR.MIDDLE, wrap=False, margins=(2, 2, 1, 1))


def add_bullets(slide, items, x, y, width, text_color, bullet_color, font_size=18, gap=12):
    current_y = y
    for item in items:
        add_circle(slide, x, current_y + 8, 8, bullet_color, line=None)
        text_h = estimate_text_height(item, FONT_BODY, font_size, width - 40, spacing=1.0) + 8
        add_textbox(slide, x + 20, current_y, width - 20, text_h, item, font_size=font_size, color=text_color, bold=True if font_size >= 19 else False)
        current_y += text_h + gap
    return current_y


def add_window_card(slide, title, body, x, y, w, h, logo_name, accent, theme):
    add_rect(slide, x, y, w, h, fill=theme["panel"], line=theme["line"], line_width=1.25)
    add_rect(slide, x, y, w, 38, fill=theme["panel2"], line=None, rounded=True)
    for idx, color in enumerate(["#ff615c", "#ffbd44", "#00ca4e"]):
        add_circle(slide, x + 14 + idx * 18, y + 10, 10, color, line=None)
    add_picture_contain(slide, LOGOS / logo_name, x + 42, y + 8, 24, 22)
    add_textbox(slide, x + 72, y + 4, w - 90, 30, title, font_size=17, color=theme["text"], bold=True, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
    add_rect(slide, x + 20, y + 50, w - 40, 44, fill=theme["accent_bar_bg"], line=None)
    add_textbox(slide, x + 30, y + 55, w - 60, 32, body, font_size=13, color=accent, bold=True, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
    line_y = y + 114
    for k in range(4):
        line_w = w - 50 if k != 3 else int((w - 50) * 0.62)
        add_rect(slide, x + 20, line_y, line_w, 14, fill=theme["line_fill"], line=None)
        line_y += 28


def add_board_tile(slide, title, body, x, y, w, h, logo_name, accent, title_font_size=18, body_font_size=14):
    add_rect(slide, x, y, w, h, fill="#ffffff", line="#d7e2ee", line_width=1.0)
    add_picture_contain(slide, LOGOS / logo_name, x + 18, y + 18, 28, 28)
    title_h = min(74, max(30, estimate_text_height(title, FONT_BODY, title_font_size, w - 74, spacing=1.0) + 4))
    add_textbox(slide, x + 56, y + 16, w - 74, title_h, title, font_size=title_font_size, color="#17314f", bold=True)
    add_rect(slide, x + 18, y + 52, 100, 5, fill=accent, line=None, rounded=False)
    body_y = y + max(70, 16 + title_h + 18)
    add_textbox(slide, x + 18, body_y, w - 36, h - (body_y - y) - 18, body, font_size=body_font_size, color="#4d6581", bold=True, spacing=1.0)


def add_vc_bg(slide):
    add_rect(slide, 0, 0, 1920, 1080, fill="#091727", line=None, rounded=False)
    tri1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, px(1460), px(-10), px(470), px(320))
    tri1.fill.solid()
    tri1.fill.fore_color.rgb = rgb("#8b5cf6")
    tri1.fill.transparency = 0.18
    tri1.line.fill.background()
    tri1.rotation = 180
    tri2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, px(-20), px(820), px(360), px(280))
    tri2.fill.solid()
    tri2.fill.fore_color.rgb = rgb("#67e8f9")
    tri2.fill.transparency = 0.14
    tri2.line.fill.background()


def add_board_bg(slide):
    add_rect(slide, 0, 0, 1920, 1080, fill="#f7f9fc", line=None, rounded=False)
    add_rect(slide, 0, 0, 1920, 12, fill="#0b6bcb", line=None, rounded=False)
    add_rect(slide, 1540, 0, 380, 1080, fill="#eef3f8", line=None, rounded=False)


def add_vc_header(slide, title, subtitle):
    add_textbox(slide, 90, 56, 1600, 154, title, font_name="Avenir Next Condensed", font_size=42, color="#f6f1e8", bold=True, spacing=0.95)
    if subtitle:
        add_textbox(slide, 90, 226, 1600, 76, subtitle, font_name="Avenir Next", font_size=18, color="#a8b5c8", bold=True, spacing=1.0)


def add_board_header(slide, title, subtitle):
    add_textbox(slide, 120, 68, 1320, 148, title, font_name="Avenir Next Condensed", font_size=36, color="#17314f", bold=True, spacing=0.95)
    add_textbox(slide, 120, 222, 1300, 72, subtitle, font_name="Avenir Next", font_size=18, color="#667d96", bold=True, spacing=1.0)


def vc_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_vc_bg(slide)
    add_rect(slide, 100, 84, 270, 42, fill="#10253f", line="#29415f", line_width=1.0)
    add_textbox(slide, 118, 88, 240, 32, "KIRKLAND AI × EMMANEIGH", font_size=14, color="#a8b5c8", bold=True)
    add_textbox(slide, 100, 200, 760, 280, "The model can reason.\nEmmaNeigh makes\nit execute.", font_name="Avenir Next Condensed", font_size=50, color="#f6f1e8", bold=True)
    add_textbox(slide, 100, 540, 770, 110, "A model-agnostic execution layer for matter administration, paralegal workflows, and document operations across the legal stack.", font_size=21, color="#a8b5c8", bold=True)
    add_rect(slide, 100, 702, 660, 104, fill="#12233d", line="#304661")
    add_textbox(slide, 130, 722, 600, 42, "Execution is the moat.", font_name="Avenir Next Condensed", font_size=26, color="#fbbf24", bold=True)
    add_textbox(slide, 130, 764, 500, 28, "Models will change. The workflow fabric compounds.", font_size=16, color="#f6f1e8", bold=True)
    add_rect(slide, 1130, 180, 640, 650, fill="#0c1828", line="#24364d")
    add_textbox(slide, 1180, 216, 320, 42, "Execution fabric", font_name="Avenir Next Condensed", font_size=28, color="#f6f1e8", bold=True)
    add_textbox(slide, 1180, 270, 520, 70, "AI interprets intent. EmmaNeigh routes and executes the work across the actual operating stack.", font_size=18, color="#a8b5c8", bold=True)
    chips = [("Harvey", 1180, 360, 110), ("Claude", 1300, 360, 110), ("GPT", 1420, 360, 104), ("Groq / Qwen", 1538, 360, 142), ("Local model", 1180, 420, 160)]
    for label, x, y, w in chips:
        add_text_chip(slide, label, x, y, w, 42, "#111b2d", "#29415f", "#f6f1e8", font_size=16)
    add_rect(slide, 1180, 500, 550, 70, fill="#17304b", line="#5ec8ff")
    add_textbox(slide, 1208, 516, 500, 36, "EmmaNeigh execution layer", font_name="Avenir Next Condensed", font_size=22, color="#f6f1e8", bold=True)
    chip_specs = [
        ("Outlook", "outlook.png", 1180, 620, 175),
        ("Adobe", "adobe.png", 1370, 620, 160),
        ("DocuSign", "docusign.png", 1540, 620, 190),
        ("iManage", "imanage.png", 1180, 700, 205),
        ("Litera", "litera_favicon.png", 1400, 700, 160),
    ]
    for label, logo, x, y, w in chip_specs:
        add_logo_chip(slide, label, logo, x, y, w, 52, "#0e1d30", "#304661", "#f6f1e8", font_size=16)
    add_textbox(slide, 100, 980, 420, 24, "Confidential discussion draft · April 2026", font_size=12, color="#7d8ea7", bold=True)


def vc_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_vc_bg(slide)
    add_vc_header(slide, "The bottleneck is not thinking. It is\nexecution across the desktop stack.", "Matter administration still requires people to open files, switch systems, and complete deterministic steps one application at a time.")
    add_textbox(slide, 110, 318, 660, 30, "A single matter can require live interaction across all of this:", font_size=16, color="#a8b5c8", bold=True)
    theme = {"panel": "#0d1b2c", "panel2": "#111f34", "line": "#304661", "accent_bar_bg": "#10253f", "line_fill": "#17304b", "text": "#f6f1e8"}
    add_window_card(slide, "Outlook", "Find the latest draft, save attachments, confirm who sent what", 110, 362, 360, 250, "outlook.png", "#38bdf8", theme)
    add_window_card(slide, "Adobe", "Split, combine, clean, and package PDFs for execution and closing", 500, 362, 360, 250, "adobe.png", "#fb7185", theme)
    add_window_card(slide, "iManage", "Save down, version up, browse versions, organize matter files", 110, 642, 360, 250, "imanage.png", "#67e8f9", theme)
    add_window_card(slide, "Litera + DocuSign", "Run redlines, create sig packets, assemble executed pages", 500, 642, 360, 250, "docusign.png", "#fb923c", theme)
    dot_y = 350
    while dot_y < 870:
        add_circle(slide, 918, dot_y, 12, "#284364", line=None)
        dot_y += 28
    add_textbox(slide, 872, 898, 150, 22, "SYSTEM BOUNDARY", font_size=11, color="#7d8ea7", bold=True)
    add_rect(slide, 980, 362, 840, 470, fill="#0b1322", line="#304863")
    add_textbox(slide, 1010, 384, 300, 22, "Representative cloud-AI constraint", font_size=13, color="#fbbf24", bold=True)
    add_logo_chip(slide, "Harvey / other reasoning model", "harvey.png", 1010, 418, 420, 40, "#171315", "#2d2a30", "#f6f1e8", font_size=15)
    add_textbox(slide, 1010, 488, 720, 140, "“I can analyze this agreement, but I cannot create a sig packet, run Litera, or save a version into iManage from your desktop.”", font_name="Avenir Next Condensed", font_size=28, color="#f6f1e8", bold=True)
    add_textbox(slide, 1010, 706, 700, 48, "The reasoning layer is in the cloud. The work still lives across local files, desktop applications, and firm systems.", font_size=18, color="#a8b5c8", bold=True)
    add_rect(slide, 90, 902, 1720, 84, fill="#0b1524", line="#2a405c")
    add_textbox(slide, 118, 928, 1600, 28, "Result: attorneys and paralegals still do high-volume operational work manually, and partners often write the time off anyway.", font_size=16, color="#fbbf24", bold=True)


def vc_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_vc_bg(slide)
    add_vc_header(slide, "EmmaNeigh unifies the AI layer and the\nsystems-integration layer.", "Think of it as a USB-C hub for legal work: one execution layer that plugs models into Outlook, Adobe, Litera, iManage, DocuSign, files, and checklists.")
    add_textbox(slide, 110, 318, 380, 30, "1. Model-agnostic reasoning layer", font_size=15, color="#67e8f9", bold=True)
    chips = [("Harvey", 110, 354, 84), ("Claude", 210, 354, 84), ("GPT", 310, 354, 88), ("Groq / Qwen", 414, 354, 150), ("Local model", 584, 354, 150)]
    centers = []
    for label, x, y, w in chips:
        add_text_chip(slide, label, x, y, w, 44, "#111b2d", "#29415f", "#f6f1e8", font_size=15)
        centers.append(x + w / 2)
    for center in centers:
        add_line(slide, center, 398, center, 438, "#67e8f9", width=2.25)
        tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, px(center - 8), px(436), px(16), px(16))
        tri.fill.solid()
        tri.fill.fore_color.rgb = rgb("#67e8f9")
        tri.line.fill.background()
        tri.rotation = 180
    add_rect(slide, 110, 458, 1700, 90, fill="#17304b", line="#5ec8ff")
    add_textbox(slide, 148, 479, 520, 36, "2. EmmaNeigh execution layer", font_name="Avenir Next Condensed", font_size=26, color="#f6f1e8", bold=True)
    add_textbox(slide, 1020, 472, 700, 40, "Translates natural-language intent into machine-readable commands and then executes them deterministically.", font_size=15, color="#a8b5c8", bold=True)
    add_textbox(slide, 110, 598, 560, 26, "3. Actual operating systems and applications", font_size=15, color="#fb923c", bold=True)
    add_text_chip(slide, "Local files", 110, 634, 170, 54, "#0d1b2c", "#304661", "#f6f1e8", font_size=15)
    chip_specs = [
        ("Outlook", "outlook.png", 310, 634, 185),
        ("Adobe", "adobe.png", 520, 634, 170),
        ("Litera", "litera_favicon.png", 715, 634, 170),
        ("DocuSign", "docusign.png", 930, 634, 205),
        ("iManage", "imanage.png", 1160, 634, 250),
    ]
    for label, logo, x, y, w in chip_specs:
        add_logo_chip(slide, label, logo, x, y, w, 54, "#0d1b2c", "#304661", "#f6f1e8", font_size=15)
    bullets = [
        "AI is used for translation and routing, not for performing the underlying deterministic work.",
        "No browser-click agent is required for core execution. The value comes from software integrations, file operations, and operating-system access.",
        "Because the layer is model agnostic, the reasoning system can change without rebuilding the workflow engine.",
    ]
    add_bullets(slide, bullets, 110, 710, 1540, "#f6f1e8", "#fbbf24", font_size=15, gap=10)


def vc_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_vc_bg(slide)
    add_vc_header(slide, "What EmmaNeigh already executes", "The product is already organized around recurring legal operations, not generic chat.")
    cards = [
        ("Email & attachments", "Search folders, determine whether a draft was received or sent, save attachments, and prepare follow-ups.", "outlook.png"),
        ("Checklist updates", "Upload a checklist, scan Outlook activity, and update comments/status based on actual matter traffic.", "outlook.png"),
        ("Punchlists", "Turn a working checklist into a cleaner punchlist format for transaction management and follow-up.", "imanage.png"),
        ("Redlines", "Run Litera comparisons and output full-document or targeted comparison sets.", "litera_favicon.png"),
        ("PDF workflows", "Split, combine, clean, and prep signature and closing PDFs.", "adobe.png"),
        ("Executed versions", "Match signed pages back into the correct agreements and rebuild executed sets at scale.", "docusign.png"),
        ("Document management", "Browse, save, organize, and version documents where the environment allows it.", "imanage.png"),
        ("File operations", "Convert files, save outputs, move deliverables, and handle repetitive desktop document actions.", "adobe.png"),
    ]
    x0, y0, w, h, gx, gy = 90, 320, 410, 170, 36, 24
    for idx, (title, body, logo) in enumerate(cards):
        row, col = divmod(idx, 4)
        x = x0 + col * (w + gx)
        y = y0 + row * (h + gy)
        add_rect(slide, x, y, w, h, fill="#0d1b2c", line="#304661")
        add_picture_contain(slide, LOGOS / logo, x + 18, y + 18, 28, 28)
        add_textbox(slide, x + 56, y + 14, w - 74, 34, title, font_name="Avenir Next Condensed", font_size=20, color="#f6f1e8", bold=True)
        add_textbox(slide, x + 18, y + 58, w - 36, h - 70, body, font_size=15, color="#a8b5c8", bold=True, spacing=1.0)
    add_rect(slide, 90, 904, 1740, 78, fill="#0c1828", line="#273d58")
    add_textbox(slide, 120, 926, 1660, 32, "Core point: the execution layer is useful even without AI. AI simply makes it accessible through natural-language prompts instead of menus and macros.", font_size=17, color="#f6f1e8", bold=True)


def vc_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_vc_bg(slide)
    add_vc_header(slide, "Why firms care: this is matter administration and\nparalegal work that leaks into lawyer time.", None)
    cols = [
        ("Matter administration", "Checklists, packet assembly, version handling, distribution tracking, and status chasing are necessary but operational.", "#fbbf24"),
        ("Paralegal workflows", "Much of the work is rules-based, repetitive, and document-centric. It does not require high-end legal reasoning to execute correctly.", "#67e8f9"),
        ("Partner write-offs", "When lawyers do these tasks anyway, the time is hard to bill, easy to discount, and distracting from higher-value work.", "#fb923c"),
    ]
    x = 90
    for title, body, accent in cols:
        add_rect(slide, x, 300, 540, 326, fill="#0d1a2c", line="#273d58")
        add_textbox(slide, x + 34, 332, 430, 64, title, font_name="Avenir Next Condensed", font_size=26, color="#f6f1e8", bold=True, spacing=0.95)
        add_rect(slide, x + 34, 396, 104, 6, fill=accent, line=None, rounded=False)
        add_textbox(slide, x + 34, 430, 450, 148, body, font_size=16, color="#a8b5c8", bold=True, spacing=1.0)
        x += 590
    add_rect(slide, 90, 680, 1740, 190, fill="#0e2137", line="#2b4668")
    add_textbox(slide, 124, 710, 360, 26, "The pitch to the firm is simple:", font_size=15, color="#67e8f9", bold=True)
    add_textbox(slide, 124, 752, 1580, 104, "EmmaNeigh does not replace legal judgment. It compresses the operational execution time around documents, signatures, redlines, checklists, and matter administration so lawyers spend less time on work clients resist paying for.", font_size=17, color="#f6f1e8", bold=True, spacing=1.0)


def vc_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_vc_bg(slide)
    add_vc_header(slide, "Why this is interesting for Kirkland AI", "EmmaNeigh is not another standalone model. It is the execution surface that can sit beneath the reasoning layer Kirkland chooses to use.")
    add_rect(slide, 90, 288, 950, 610, fill="#0c1728", line="#273d58")
    add_textbox(slide, 130, 320, 300, 30, "Reasoning system", font_size=15, color="#67e8f9", bold=True)
    add_logo_chip(slide, "Harvey", "harvey.png", 130, 360, 190, 54, "#171315", "#2d2a30", "#f6f1e8", font_size=16)
    for idx, label in enumerate(["Other frontier model", "Open model", "Future internal model"]):
        add_text_chip(slide, label, 420, 360 + idx * 70, 310, 46, "#121d2f", "#324964", "#f6f1e8", font_size=15)
    add_line(slide, 310, 414, 310, 550, "#67e8f9", width=2)
    add_line(slide, 310, 480, 570, 480, "#67e8f9", width=2)
    add_rect(slide, 130, 560, 850, 74, fill="#17304b", line="#5ec8ff")
    add_textbox(slide, 160, 576, 420, 30, "EmmaNeigh execution layer", font_name="Avenir Next Condensed", font_size=22, color="#f6f1e8", bold=True)
    add_textbox(slide, 160, 658, 240, 24, "Firm operating stack", font_size=14, color="#fb923c", bold=True)
    xx = 130
    for label, logo, w in [("Outlook", "outlook.png", 180), ("Adobe", "adobe.png", 170), ("Litera", "litera_favicon.png", 170), ("iManage", "imanage.png", 220)]:
        add_logo_chip(slide, label, logo, xx, 696, w, 52, "#0e1d30", "#304661", "#f6f1e8", font_size=15)
        xx += w + 20
    add_rect(slide, 1120, 288, 710, 610, fill="#0c1728", line="#273d58")
    add_textbox(slide, 1160, 324, 520, 94, "What Kirkland gets if this\nworks", font_name="Avenir Next Condensed", font_size=32, color="#f6f1e8", bold=True)
    bullets = [
        "A model-agnostic control plane that can work with Harvey today and a different reasoning system tomorrow.",
        "A way to extend the AI discussion from drafting and analysis into concrete execution across the existing desktop and DMS stack.",
        "A path to operational leverage without requiring lawyers to manually bridge the gap between cloud AI and local applications.",
        "No rip-and-replace: the product is designed to sit on top of the tools the firm already uses.",
    ]
    add_bullets(slide, bullets, 1160, 446, 600, "#f6f1e8", "#fbbf24", font_size=17, gap=12)


def vc_slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_vc_bg(slide)
    add_vc_header(slide, "Two asks would materially improve the\nproduct in a Kirkland environment.", "These are not cosmetic asks. They unlock the exact execution depth that the current stack otherwise blocks.")
    add_rect(slide, 90, 320, 810, 450, fill="#0c1828", line="#273d58")
    add_textbox(slide, 130, 360, 640, 84, "1. Harvey bearer /\nauthentication token", font_name="Avenir Next Condensed", font_size=30, color="#f6f1e8", bold=True, spacing=0.95)
    add_textbox(slide, 130, 452, 660, 118, "That would allow EmmaNeigh to plug into a reasoning layer with stronger legal performance than a free foundation model, while preserving the model-agnostic execution architecture.", font_size=16, color="#a8b5c8", bold=True, spacing=1.0)
    add_logo_chip(slide, "Harvey", "harvey.png", 130, 590, 170, 46, "#171315", "#2d2a30", "#f6f1e8", font_size=15)
    add_rect(slide, 130, 650, 590, 92, fill="#17304b", line="#5ec8ff")
    add_textbox(slide, 160, 674, 180, 24, "Why it matters", font_size=14, color="#67e8f9", bold=True)
    add_textbox(slide, 160, 706, 520, 42, "Better planning, better task routing, better explanation quality — without changing the execution engine underneath.", font_size=14, color="#f6f1e8", bold=True, spacing=1.0)
    add_rect(slide, 1020, 320, 810, 450, fill="#0c1828", line="#273d58")
    add_textbox(slide, 1060, 360, 640, 84, "2. Unrestricted iManage API\nsurface", font_name="Avenir Next Condensed", font_size=30, color="#f6f1e8", bold=True, spacing=0.95)
    add_textbox(slide, 1060, 452, 660, 118, "Right now the accessible interface is constrained. That limits browse, version-aware workflows, and redlining across versions. A richer surface would let the product operate far more seamlessly inside the document system.", font_size=16, color="#a8b5c8", bold=True, spacing=1.0)
    add_logo_chip(slide, "iManage", "imanage.png", 1060, 590, 210, 46, "#0e1d30", "#304661", "#f6f1e8", font_size=15)
    add_rect(slide, 1060, 650, 670, 92, fill="#17304b", line="#5ec8ff")
    add_textbox(slide, 1090, 674, 180, 24, "Why it matters", font_size=14, color="#67e8f9", bold=True)
    add_textbox(slide, 1090, 706, 600, 42, "The product gets closer to full browse / save / version / compare workflows instead of stopping at partial desktop integration.", font_size=14, color="#f6f1e8", bold=True, spacing=1.0)
    add_textbox(slide, 90, 918, 1700, 40, "If Kirkland AI wants to evaluate whether this can become a real operating layer rather than just another demo, these are the leverage points.", font_size=15, color="#fbbf24", bold=True, spacing=1.0)


def vc_slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_vc_bg(slide)
    add_vc_header(slide, "A sensible next step is a narrow pilot, not a\nbroad rollout.", "The product is strongest when evaluated against a handful of real transaction workflows with the right system access.")
    add_rect(slide, 90, 320, 1090, 540, fill="#0c1828", line="#273d58")
    add_textbox(slide, 130, 356, 520, 42, "Suggested pilot structure", font_name="Avenir Next Condensed", font_size=28, color="#f6f1e8", bold=True)
    steps = [
        ("Week 1", "Connect model layer + target integrations", "#67e8f9"),
        ("Week 2", "Validate 3 workflows end-to-end on sample matters", "#fbbf24"),
        ("Weeks 3–5", "Run a small user group, measure cycle-time reduction, and identify the permissions that unlock the next tier of value.", "#fb923c"),
    ]
    yy = 450
    for idx, (head, sub, color) in enumerate(steps):
        add_circle(slide, 150, yy + 2, 18, color, line=None)
        if idx < len(steps) - 1:
            add_line(slide, 159, yy + 20, 159, yy + 82, color, width=3)
        add_textbox(slide, 210, yy - 8, 220, 32, head, font_name="Avenir Next Condensed", font_size=22, color="#f6f1e8", bold=True)
        add_textbox(slide, 210, yy + 34, 790, 58, sub, font_size=15, color="#a8b5c8", bold=True, spacing=1.0)
        yy += 128
    add_rect(slide, 1240, 320, 590, 540, fill="#0c1828", line="#273d58")
    add_textbox(slide, 1280, 356, 450, 90, "What success would look\nlike", font_name="Avenir Next Condensed", font_size=26, color="#f6f1e8", bold=True, spacing=0.95)
    outcomes = [
        ("Faster checklist maintenance", "Comments and status updated from actual matter activity"),
        ("Less manual packet work", "Signature and executed-version assembly compressed materially"),
        ("Lower write-off pressure", "Less lawyer time spent on operational execution"),
        ("Clear integration roadmap", "Know exactly which permissions unlock the next tier of value"),
    ]
    yy = 450
    for head, sub in outcomes:
        add_rect(slide, 1280, yy, 500, 94, fill="#17304b", line="#304661")
        add_textbox(slide, 1308, yy + 12, 430, 30, head, font_name="Avenir Next Condensed", font_size=21, color="#f6f1e8", bold=True)
        add_textbox(slide, 1308, yy + 46, 430, 34, sub, font_size=13, color="#a8b5c8", bold=True, spacing=1.0)
        yy += 106


def board_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_board_bg(slide)
    add_rect(slide, 120, 70, 240, 36, fill="#eef3f8", line="#d8e1ec")
    add_textbox(slide, 138, 76, 210, 28, "KIRKLAND AI × EMMANEIGH", font_size=13, color="#60758d", bold=True)
    add_textbox(slide, 120, 150, 760, 200, "The model can reason.\nEmmaNeigh makes it execute.", font_name="Avenir Next Condensed", font_size=42, color="#17314f", bold=True, spacing=0.94)
    add_textbox(slide, 120, 368, 780, 96, "A model-agnostic execution layer for matter administration, paralegal workflows, and document operations across the legal stack.", font_size=20, color="#60758d", bold=True, spacing=1.0)
    add_rect(slide, 1035, 118, 745, 724, fill="#ffffff", line="#d8e1ec")
    add_textbox(slide, 1090, 156, 320, 46, "Execution layer", font_name="Avenir Next Condensed", font_size=24, color="#17314f", bold=True)
    add_textbox(slide, 1090, 220, 620, 112, "EmmaNeigh sits between the model and the operating stack. It turns intent into executable commands and routes the work across the tools lawyers already use.", font_size=18, color="#39506b", bold=True, spacing=1.06)
    add_logo_chip(slide, "Harvey", "harvey.png", 1090, 366, 180, 48, "#ffffff", "#d8e1ec", "#17314f", font_size=15)
    add_text_chip(slide, "Firm-approved model", 1290, 366, 248, 48, "#ffffff", "#d8e1ec", "#17314f", font_size=14)
    add_text_chip(slide, "Open model", 1556, 366, 150, 48, "#ffffff", "#d8e1ec", "#17314f", font_size=14)
    add_text_chip(slide, "Future internal model", 1090, 424, 272, 48, "#ffffff", "#d8e1ec", "#17314f", font_size=13)
    add_text_chip(slide, "Managed remote model", 1382, 424, 270, 48, "#ffffff", "#d8e1ec", "#17314f", font_size=13)
    add_rect(slide, 1090, 534, 620, 82, fill="#edf5ff", line="#90bff0")
    add_textbox(slide, 1124, 550, 460, 36, "EmmaNeigh execution fabric", font_name="Avenir Next Condensed", font_size=20, color="#17314f", bold=True)
    chips = [("Outlook", "outlook.png", 1090, 670, 190), ("Word", "word.png", 1295, 670, 150), ("Adobe", "adobe.png", 1460, 670, 175), ("DocuSign", "docusign.png", 1090, 730, 200), ("Litera", "litera_favicon.png", 1305, 730, 150), ("iManage", "imanage.png", 1470, 730, 215)]
    for label, logo, x, y, w in chips:
        add_logo_chip(slide, label, logo, x, y, w, 46, "#ffffff", "#d8e1ec", "#17314f", font_size=14)


def board_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_board_bg(slide)
    add_board_header(slide, "The bottleneck is execution across the\ndesktop stack.", "Cloud AI can answer questions. The actual work still lives across local files, desktop applications, and firm systems.")
    add_textbox(slide, 120, 300, 280, 28, "Representative matter stack", font_size=13, color="#60758d", bold=True)
    tiles = [
        ("Outlook", "Search matter traffic, save attachments, confirm who sent what", "outlook.png", "#0b6bcb"),
        ("Litera", "Run redlines and produce comparison output", "litera_favicon.png", "#d97706"),
        ("iManage", "Browse, save, version, and organize matter files", "imanage.png", "#12806a"),
        ("Adobe", "Clean PDFs, remove pages, combine sets, and prep signature materials", "adobe.png", "#c2410c"),
        ("DocuSign", "Circulate signature pages and rebuild executed versions", "docusign.png", "#7c3aed"),
        ("Word", "Prepare clean documents and convert deliverables", "word.png", "#6b8df8"),
    ]
    x0, y0 = 120, 336
    for idx, (title, body, logo, accent) in enumerate(tiles):
        row, col = divmod(idx, 3)
        add_board_tile(slide, title, body, x0 + col * 290, y0 + row * 184, 260, 168, logo, accent, title_font_size=17, body_font_size=13)
    add_rect(slide, 980, 326, 820, 430, fill="#f3f5f8", line="#d8e1ec")
    add_picture_contain(slide, LOGOS / "harvey.png", 1006, 344, 26, 26)
    add_textbox(slide, 1042, 344, 120, 24, "Harvey", font_size=13, color="#60758d", bold=True)
    add_rect(slide, 1010, 378, 760, 62, fill="#ffffff", line="#d8e1ec")
    add_textbox(slide, 1036, 392, 720, 34, "Can you take this PDF and remove the signature pages to create a signature packet?", font_size=14, color="#334155", bold=True, spacing=1.0)
    add_textbox(slide, 1010, 460, 320, 24, "Copy      Save prompt      Edit query", font_size=10, color="#8b9bb0", bold=True)
    add_textbox(slide, 1010, 490, 220, 26, "Finished in 3 steps", font_size=13, color="#1f2937", bold=True)
    add_textbox(slide, 1010, 532, 710, 138, "I appreciate the request, but I'm not able to create, edit, or modify PDF files or any other file types. I can only read and analyze the content of uploaded documents, not produce new versions of them.", font_size=15, color="#334155", bold=True, spacing=1.0)
    add_rect(slide, 980, 780, 820, 118, fill="#fff7e8", line="#edd6ab")
    add_textbox(slide, 1010, 806, 760, 66, "Key point: the model is not refusing because the task is unintelligent. It is refusing because it does not have the operating-system and application-level execution surface.", font_size=15, color="#b7791f", bold=True, spacing=1.0)


def board_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_board_bg(slide)
    add_board_header(slide, "EmmaNeigh unifies reasoning and\nexecution.", "AI determines what should happen. EmmaNeigh determines how to execute it across connected software and files.")
    add_rect(slide, 120, 314, 1680, 606, fill="#ffffff", line="#d8e1ec")
    add_textbox(slide, 170, 350, 320, 44, "1. Reasoning layer", font_name="Avenir Next Condensed", font_size=25, color="#17314f", bold=True)
    add_textbox(slide, 170, 400, 360, 70, "Harvey or another model interprets the request, identifies the workflow, and fills in the parameters.", font_size=16, color="#39506b", bold=True)
    add_logo_chip(slide, "Harvey", "harvey.png", 170, 492, 150, 40, "#ffffff", "#d8e1ec", "#17314f", font_size=15)
    add_text_chip(slide, "Other model", 170, 544, 180, 40, "#ffffff", "#d8e1ec", "#17314f", font_size=15)
    add_bullets(slide, [
        "Execution does not require AI; it requires connectors, file operations, and system access.",
        "EmmaNeigh stays model agnostic as reasoning models improve. AI simply makes the workflows easier to invoke.",
    ], 170, 592, 470, "#17314f", "#b7791f", font_size=15, gap=10)
    add_textbox(slide, 760, 350, 460, 44, "2. EmmaNeigh execution layer", font_name="Avenir Next Condensed", font_size=25, color="#17314f", bold=True)
    add_rect(slide, 760, 400, 980, 96, fill="#edf5ff", line="#90bff0")
    add_textbox(slide, 800, 424, 860, 44, "Turns natural-language intent into machine-readable commands and deterministic workflows.", font_name="Avenir Next Condensed", font_size=22, color="#163962", bold=True)
    add_textbox(slide, 760, 534, 700, 44, "3. Operating systems and applications", font_name="Avenir Next Condensed", font_size=25, color="#17314f", bold=True)
    row1 = [("Local files", "adobe.png", 190), ("Outlook", "outlook.png", 180), ("Word", "word.png", 150), ("Adobe", "adobe.png", 165), ("Litera", "litera_favicon.png", 150)]
    row2 = [("DocuSign", "docusign.png", 190), ("iManage", "imanage.png", 220)]
    xx = 760
    for label, logo, w in row1:
        add_logo_chip(slide, label, logo, xx, 594, w, 42, "#ffffff", "#d8e1ec", "#17314f", font_size=14)
        xx += w + 14
    xx = 760
    for label, logo, w in row2:
        add_logo_chip(slide, label, logo, xx, 648, w, 42, "#ffffff", "#d8e1ec", "#17314f", font_size=14)
        xx += w + 16


def board_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_board_bg(slide)
    add_board_header(slide, "What EmmaNeigh already executes", "The workflow surface is already practical and operational, not theoretical.")
    cards = [
        ("Checklist updates", "Upload a checklist, scan matter activity, and update comments and status.", "outlook.png", "#0b6bcb"),
        ("Punchlists", "Turn a working checklist into a cleaner matter-management punchlist.", "imanage.png", "#6b8df8"),
        ("Redlines", "Run Litera comparisons and generate comparison outputs.", "litera_favicon.png", "#d97706"),
        ("PDF workflows", "Remove signature pages, combine sets, and prep execution-ready PDF materials.", "adobe.png", "#c2410c"),
        ("Executed versions", "Match signed pages back into agreements and rebuild executed sets.", "docusign.png", "#7c3aed"),
        ("Document movement", "Move, convert, save, and organize deliverables across files and connected systems.", "imanage.png", "#12806a"),
    ]
    x0, y0, w, h = 120, 316, 530, 182
    for idx, (title, body, logo, accent) in enumerate(cards):
        row, col = divmod(idx, 2)
        add_board_tile(slide, title, body, x0 + col * 554, y0 + row * 206, w, h, logo, accent)
    add_rect(slide, 1248, 316, 552, 560, fill="#ffffff", line="#d8e1ec")
    add_textbox(slide, 1290, 350, 300, 44, "Operating principle", font_name="Avenir Next Condensed", font_size=25, color="#17314f", bold=True)
    add_bullets(slide, [
        "These workflows can be invoked manually or via AI routing.",
        "The real moat is not the model. It is the execution layer across the legal stack.",
        "The same workflow engine can work beneath Harvey today and another reasoning layer later.",
    ], 1290, 442, 450, "#17314f", "#0b6bcb", font_size=16, gap=18)


def board_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_board_bg(slide)
    add_board_header(slide, "Why firms care", "This is matter administration and paralegal work that often leaks into lawyer time.")
    cols = [
        ("Matter administration", "Checklists, packet assembly, version handling, distribution tracking, and status chasing are necessary but operational.", "#b7791f"),
        ("Paralegal workflows", "Much of the work is rules-based, repetitive, and document-centric. It does not require high-end legal reasoning to execute correctly.", "#0b6bcb"),
        ("Partner write-offs", "When lawyers do these tasks anyway, the time is hard to bill, easy to discount, and distracting from higher-value work.", "#d97706"),
    ]
    x = 120
    for title, body, accent in cols:
        add_rect(slide, x, 320, 530, 356, fill="#ffffff", line="#d8e1ec")
        add_textbox(slide, x + 34, 356, 390, 68, title, font_name="Avenir Next Condensed", font_size=24, color="#17314f", bold=True, spacing=0.95)
        add_rect(slide, x + 34, 414, 130, 6, fill=accent, line=None, rounded=False)
        add_textbox(slide, x + 34, 446, 450, 164, body, font_size=15, color="#39506b", bold=True, spacing=1.0)
        x += 570
    add_rect(slide, 120, 734, 1680, 176, fill="#ffffff", line="#d8e1ec")
    add_textbox(slide, 160, 762, 320, 28, "Boardroom framing", font_size=15, color="#0b6bcb", bold=True)
    add_textbox(slide, 160, 808, 1540, 94, "EmmaNeigh does not replace legal judgment. It compresses the operational execution time around documents, signatures, redlines, checklists, PDFs, and matter administration so lawyers spend less time on work clients resist paying for.", font_name="Avenir Next Condensed", font_size=21, color="#17314f", bold=True, spacing=0.95)


def board_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_board_bg(slide)
    add_board_header(slide, "Why this is interesting for Kirkland AI", "EmmaNeigh is not another model discussion. It is the execution surface that can sit beneath the firm’s preferred reasoning layer.")
    add_rect(slide, 120, 314, 820, 586, fill="#ffffff", line="#d8e1ec")
    add_textbox(slide, 160, 350, 420, 50, "Kirkland AI implication", font_name="Avenir Next Condensed", font_size=25, color="#17314f", bold=True)
    add_textbox(slide, 160, 400, 700, 88, "If Kirkland wants to explore Harvey or another reasoning system beyond drafting and analysis, EmmaNeigh is the layer that can make those systems operational across the existing desktop and DMS environment.", font_size=16, color="#39506b", bold=True)
    add_textbox(slide, 160, 512, 260, 28, "Illustrative stack", font_size=13, color="#60758d", bold=True)
    chips = [("Harvey", "harvey.png", 160, 550, 150), ("Outlook", "outlook.png", 345, 550, 180), ("Adobe", "adobe.png", 560, 550, 165), ("Litera", "litera_favicon.png", 160, 614, 150), ("DocuSign", "docusign.png", 345, 614, 195), ("iManage", "imanage.png", 560, 614, 220)]
    for label, logo, x, y, w in chips:
        add_logo_chip(slide, label, logo, x, y, w, 42, "#ffffff", "#d8e1ec", "#17314f", font_size=14)
    add_rect(slide, 1000, 314, 800, 586, fill="#ffffff", line="#d8e1ec")
    add_textbox(slide, 1040, 350, 520, 44, "What Kirkland gets if this works", font_name="Avenir Next Condensed", font_size=25, color="#17314f", bold=True)
    add_bullets(slide, [
        "A model-agnostic control plane that can work with Harvey today and a different reasoning system tomorrow.",
        "A path from AI analysis into concrete execution across the existing desktop and DMS stack.",
        "Operational leverage without requiring lawyers to manually bridge the gap between cloud AI and local systems.",
        "No rip-and-replace: the product sits on top of the tools the firm already uses.",
    ], 1040, 424, 680, "#17314f", "#b7791f", font_size=17, gap=16)


def board_slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_board_bg(slide)
    add_board_header(slide, "Two inputs would materially improve the product at Kirkland.", "These are not cosmetic asks. They are the two leverage points that would make the product materially more seamless.")
    add_rect(slide, 120, 320, 790, 520, fill="#ffffff", line="#d8e1ec")
    add_textbox(slide, 160, 360, 620, 62, "1. Harvey bearer / authentication token", font_name="Avenir Next Condensed", font_size=25, color="#17314f", bold=True, spacing=0.95)
    add_textbox(slide, 160, 432, 690, 104, "This would allow EmmaNeigh to plug into a reasoning layer with stronger legal performance than a free foundation model while preserving the model-agnostic execution architecture.", font_size=15, color="#39506b", bold=True, spacing=1.0)
    add_logo_chip(slide, "Harvey", "harvey.png", 160, 550, 170, 44, "#ffffff", "#d8e1ec", "#17314f", font_size=15)
    add_rect(slide, 160, 620, 700, 120, fill="#edf5ff", line="#90bff0")
    add_textbox(slide, 190, 646, 180, 28, "Why it matters", font_size=14, color="#0b6bcb", bold=True)
    add_textbox(slide, 190, 682, 620, 52, "Better planning, better task routing, and better explanation quality without changing the execution engine underneath.", font_size=15, color="#17314f", bold=True, spacing=1.0)
    add_rect(slide, 1010, 320, 790, 520, fill="#ffffff", line="#d8e1ec")
    add_textbox(slide, 1050, 360, 620, 62, "2. Unrestricted iManage API surface", font_name="Avenir Next Condensed", font_size=25, color="#17314f", bold=True, spacing=0.95)
    add_textbox(slide, 1050, 432, 690, 104, "Right now the accessible interface is constrained. That limits browse, version-aware workflows, and redlining across versions. A richer surface would let the product operate far more seamlessly inside the document system.", font_size=15, color="#39506b", bold=True, spacing=1.0)
    add_logo_chip(slide, "iManage", "imanage.png", 1050, 550, 210, 44, "#ffffff", "#d8e1ec", "#17314f", font_size=15)
    add_rect(slide, 1050, 620, 700, 120, fill="#edf5ff", line="#90bff0")
    add_textbox(slide, 1080, 646, 180, 28, "Why it matters", font_size=14, color="#0b6bcb", bold=True)
    add_textbox(slide, 1080, 682, 620, 52, "The product gets closer to full browse, save, version, and compare workflows instead of stopping at partial desktop integration.", font_size=15, color="#17314f", bold=True, spacing=1.0)
    add_textbox(slide, 120, 910, 1620, 42, "If the goal is to evaluate whether EmmaNeigh can become a real operating layer rather than just another demo, these are the two inputs that matter most.", font_size=14, color="#b7791f", bold=True, spacing=1.0)


def board_slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_board_bg(slide)
    add_board_header(slide, "A sensible next step is a narrow pilot, not a broad rollout.", "The strongest test is a focused pilot against a handful of real transaction workflows with the right system access.")
    add_rect(slide, 120, 320, 940, 580, fill="#ffffff", line="#d8e1ec")
    add_textbox(slide, 160, 352, 420, 44, "Suggested pilot structure", font_name="Avenir Next Condensed", font_size=25, color="#17314f", bold=True)
    steps = [
        ("Week 1", "Connect the reasoning layer and target integrations", "#0b6bcb"),
        ("Week 2", "Validate three workflows end-to-end on sample matters", "#b7791f"),
        ("Weeks 3–5", "Run a small user group, measure cycle-time reduction, and identify which permissions unlock the next tier.", "#d97706"),
    ]
    yy = 450
    for idx, (head, sub, color) in enumerate(steps):
        add_circle(slide, 180, yy + 4, 18, color, line=None)
        if idx < len(steps) - 1:
            add_line(slide, 189, yy + 22, 189, yy + 82, color, width=3)
        add_textbox(slide, 240, yy - 8, 200, 38, head, font_name="Avenir Next Condensed", font_size=22, color="#17314f", bold=True)
        add_textbox(slide, 240, yy + 34, 700, 56, sub, font_size=14, color="#39506b", bold=True, spacing=1.0)
        yy += 128
    add_rect(slide, 1120, 320, 680, 580, fill="#ffffff", line="#d8e1ec")
    add_textbox(slide, 1160, 352, 520, 50, "What success would look like", font_name="Avenir Next Condensed", font_size=24, color="#17314f", bold=True)
    outcomes = [
        ("Faster checklist maintenance", "Comments and status updated from actual matter activity"),
        ("Less manual packet work", "Signature and executed-version assembly compressed materially"),
        ("Lower write-off pressure", "Less lawyer time spent on operational execution"),
        ("Clear integration roadmap", "Know exactly which permissions unlock the next tier of value"),
    ]
    yy = 430
    for head, sub in outcomes:
        add_rect(slide, 1160, yy, 560, 104, fill="#f7fbff", line="#cfe0f4")
        add_textbox(slide, 1190, yy + 14, 500, 34, head, font_name="Avenir Next Condensed", font_size=20, color="#17314f", bold=True)
        add_textbox(slide, 1190, yy + 48, 500, 40, sub, font_size=13, color="#39506b", bold=True, spacing=1.0)
        yy += 118


def build_deck(output_path, slide_fns):
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    for fn in slide_fns:
        fn(prs)
    prs.save(str(output_path))


def main():
    build_deck(DOCS / "KIRKLAND_AI_PITCH_DECK_APRIL_2026.pptx", [
        vc_slide_1,
        vc_slide_2,
        vc_slide_3,
        vc_slide_4,
        vc_slide_5,
        vc_slide_6,
        vc_slide_7,
        vc_slide_8,
    ])
    build_deck(DOCS / "KIRKLAND_AI_PITCH_DECK_APRIL_2026_BOARDROOM.pptx", [
        board_slide_1,
        board_slide_2,
        board_slide_3,
        board_slide_4,
        board_slide_5,
        board_slide_6,
        board_slide_7,
        board_slide_8,
    ])
    print(DOCS / "KIRKLAND_AI_PITCH_DECK_APRIL_2026.pptx")
    print(DOCS / "KIRKLAND_AI_PITCH_DECK_APRIL_2026_BOARDROOM.pptx")


if __name__ == "__main__":
    main()
